from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.util import Pt

INPUT_FILE = "PPT_Template.pptx"
OUTPUT_FILE = "AI_Viewers_Hackathon_Presentation.pptx"
FALLBACK_OUTPUT_FILE = "AI_Viewers_Hackathon_Presentation_v2.pptx"

prs = Presentation(INPUT_FILE)


def set_shape_text(slide, shape_index, text):
    slide.shapes[shape_index].text = text


def add_body_text(slide, text, left, top, width, height, font_size=18, bold=False):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.font.size = Pt(font_size)
    paragraph.font.bold = bold
    paragraph.font.color.rgb = RGBColor(34, 34, 34)
    return textbox


def add_flow_box(slide, text, left, top, width, height, fill_rgb):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill_rgb)
    shape.line.color.rgb = RGBColor(255, 255, 255)
    shape.text_frame.text = text
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size = Pt(14)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)
    return shape


# Slide 2: project title and team details
slide = prs.slides[1]
set_shape_text(slide, 2, "AI Interview Platform")
set_shape_text(slide, 5, "Team Name:  | Team ID: ")
set_shape_text(slide, 9, "Team Members:\n1.\n2.")


# Slide 3: Problem Statement
slide = prs.slides[2]
add_body_text(
    slide,
    "Current Hiring Problems",
    left=1180000,
    top=1380000,
    width=3500000,
    height=350000,
    font_size=20,
    bold=True,
)
add_body_text(
    slide,
    (
        "- Manual interviews take too much recruiter time\n"
        "- Different interviewers create inconsistent evaluation\n"
        "- Scheduling slows the shortlist process\n"
        "- High-volume recruitment is difficult to scale\n"
        "- Feedback and ranking are delayed for decision makers"
    ),
    left=1180000,
    top=1760000,
    width=9000000,
    height=2500000,
    font_size=18,
)
add_body_text(
    slide,
    "Impact: slower hiring pipeline, higher screening effort, and uneven candidate assessment.",
    left=1180000,
    top=4400000,
    width=9000000,
    height=500000,
    font_size=18,
    bold=True,
)


# Slide 4: Proposed Solution / Idea with flowchart
slide = prs.slides[3]
add_body_text(
    slide,
    "Interview Flow Chart",
    left=1200000,
    top=1250000,
    width=2800000,
    height=300000,
    font_size=19,
    bold=True,
)
box_top = 1850000
box_width = 1500000
box_height = 700000
gap = 260000
start_left = 950000
flow_steps = [
    ("HR Login", (34, 87, 122)),
    ("Create Session", (46, 125, 50)),
    ("Share Link", (255, 143, 0)),
    ("AI Interview", (142, 36, 170)),
    ("Score & Rank", (198, 40, 40)),
]
for index, (label, color) in enumerate(flow_steps):
    left = start_left + index * (box_width + gap)
    add_flow_box(slide, label, left, box_top, box_width, box_height, color)
    if index < len(flow_steps) - 1:
        add_body_text(
            slide,
            "->",
            left + box_width + 70000,
            box_top + 170000,
            180000,
            200000,
            font_size=26,
            bold=True,
        )

add_body_text(
    slide,
    (
        "Solution Summary: The platform lets recruiters train an AI avatar, build question sets, "
        "send interview links, record candidate answers, and automatically generate evaluation-ready results."
    ),
    left=1180000,
    top=3150000,
    width=9300000,
    height=1400000,
    font_size=18,
)


# Slide 5: Innovation & Uniqueness with algorithm names
slide = prs.slides[4]
add_body_text(
    slide,
    "Core Algorithms Used",
    left=1180000,
    top=1300000,
    width=3400000,
    height=320000,
    font_size=19,
    bold=True,
)
add_body_text(
    slide,
    (
        "1. Keyword Matching Algorithm: compares answer text with expected role keywords.\n"
        "2. Relevance Scoring Algorithm: measures overlap between question terms and candidate response.\n"
        "3. Confidence Scoring Algorithm: evaluates confidence phrases and response length.\n"
        "4. Weighted Overall Score Algorithm: combines relevance, accuracy, confidence, and keyword score.\n"
        "5. Candidate Ranking Algorithm: sorts candidates by overall score for recruiter review."
    ),
    left=1180000,
    top=1700000,
    width=9300000,
    height=2200000,
    font_size=17,
)
add_body_text(
    slide,
    (
        "Uniqueness: the system joins avatar-led interviewing, response recording, live transcription, "
        "and rule-based AI scoring in one workflow instead of separate HR tools."
    ),
    left=1180000,
    top=4250000,
    width=9300000,
    height=1000000,
    font_size=18,
)


# Slide 6: Technology Used and modules
slide = prs.slides[5]
add_body_text(
    slide,
    "Tech Stack",
    left=1180000,
    top=1260000,
    width=1800000,
    height=300000,
    font_size=19,
    bold=True,
)
add_body_text(
    slide,
    (
        "React 18, Vite, Tailwind CSS, React Router DOM, Recharts, face-api.js,\n"
        "Web Speech API, MediaRecorder API, Lucide React, Context API"
    ),
    left=1180000,
    top=1620000,
    width=9300000,
    height=900000,
    font_size=17,
)
add_body_text(
    slide,
    "Project Modules Used",
    left=1180000,
    top=2700000,
    width=2800000,
    height=300000,
    font_size=19,
    bold=True,
)
add_body_text(
    slide,
    (
        "Modules: AuthContext, InterviewContext, AIAvatar, VideoRecorder, SpeechToText, Timer,\n"
        "CandidateRegistration, InterviewRoom, Dashboard, Sessions, CreateSession, Candidates,\n"
        "aiAnalysis.js, mediaUtils.js, helpers.js"
    ),
    left=1180000,
    top=3070000,
    width=9300000,
    height=1700000,
    font_size=17,
)


# Slide 7: Features and functional modules
slide = prs.slides[6]
add_body_text(
    slide,
    "Functional Modules",
    left=1180000,
    top=1280000,
    width=2600000,
    height=300000,
    font_size=19,
    bold=True,
)
add_body_text(
    slide,
    (
        "- Authentication Module: user login and role handling\n"
        "- Session Management Module: create, update, publish, and close interviews\n"
        "- Candidate Module: registration, interview access, and completion tracking\n"
        "- Recording Module: video and audio capture using MediaRecorder\n"
        "- Transcription Module: speech recognition for answer text\n"
        "- AI Analysis Module: keyword, relevance, confidence, and ranking logic\n"
        "- Dashboard Module: candidate list, scores, and recruiter insights"
    ),
    left=1180000,
    top=1650000,
    width=9300000,
    height=2800000,
    font_size=17,
)


# Slide 8: Implementation plan plus algorithm pipeline
slide = prs.slides[7]
add_body_text(
    slide,
    "Implementation + Algorithm Pipeline",
    left=1180000,
    top=1250000,
    width=4300000,
    height=320000,
    font_size=19,
    bold=True,
)
add_body_text(
    slide,
    (
        "Phase 1: Authentication and interview session module\n"
        "Phase 2: Avatar training and media capture module\n"
        "Phase 3: Speech-to-text and answer storage module\n"
        "Phase 4: Analysis pipeline using Keyword Match, Relevance Score, Confidence Score,\n"
        "         Weighted Overall Score, and Candidate Ranking algorithms\n"
        "Phase 5: Dashboard analytics, optimization, and production deployment"
    ),
    left=1180000,
    top=1680000,
    width=9300000,
    height=2200000,
    font_size=17,
)
add_body_text(
    slide,
    "Pipeline: Input Answer -> Transcription -> Feature Extraction -> Score Calculation -> Ranking -> Dashboard Output",
    left=1180000,
    top=4200000,
    width=9300000,
    height=900000,
    font_size=17,
    bold=True,
)

try:
    prs.save(OUTPUT_FILE)
    print(f"Created {OUTPUT_FILE}")
except PermissionError:
    fallback_path = Path(FALLBACK_OUTPUT_FILE)
    prs.save(fallback_path)
    print(f"Created {fallback_path}")
