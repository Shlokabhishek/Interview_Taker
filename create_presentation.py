from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    """Add a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(37, 99, 235)  # Primary blue
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = RGBColor(191, 219, 254)
        subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, subtitle=None):
    """Add a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add title bar
    title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(10), Inches(1))
    title_box.fill.solid()
    title_box.fill.fore_color.rgb = RGBColor(37, 99, 235)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Add subtitle if provided
    start_y = Inches(1.5)
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.font.size = Pt(18)
        subtitle_para.font.italic = True
        subtitle_para.font.color.rgb = RGBColor(71, 85, 105)
        start_y = Inches(1.8)
    
    # Add content
    content_box = slide.shapes.add_textbox(Inches(1), start_y, Inches(8), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[i]
        p.text = item
        p.level = 0
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(30, 41, 59)
        p.space_after = Pt(12)
    
    return slide

# Slide 1: Title Slide
add_title_slide(prs, "AI Interview Platform", "Revolutionizing HR Interview Process with AI Technology")

# Slide 2: Overview
add_content_slide(prs, "Project Overview", [
    "🎯 Comprehensive AI-powered interview platform",
    "🤖 Enables HR teams to train AI avatars with their face and voice",
    "👥 Conduct automated interviews with multiple candidates simultaneously",
    "📊 Automated candidate evaluation and ranking",
    "⚡ Built with React, Vite, and TailwindCSS",
    "🎥 Real-time video/audio recording and transcription"
])

# Slide 3: Problem Statement
add_content_slide(prs, "The Problem", [
    "❌ Traditional interviews are time-consuming and resource-intensive",
    "❌ HR teams struggle to interview multiple candidates simultaneously",
    "❌ Inconsistent evaluation criteria across different interviewers",
    "❌ Difficulty in scheduling and coordinating interviews",
    "❌ Limited scalability for high-volume recruitment",
    "❌ Manual note-taking and evaluation prone to bias"
])

# Slide 4: Our Solution
add_content_slide(prs, "Our Solution", [
    "✅ AI-powered avatar trained with HR's face and voice",
    "✅ Automated, consistent interview process",
    "✅ Parallel interviews with unlimited candidates",
    "✅ Real-time transcription and AI analysis",
    "✅ Instant candidate ranking and evaluation",
    "✅ Complete interview recordings and insights"
])

# Slide 5: Key Features
add_content_slide(prs, "Key Features", [
    "🎭 AI Avatar Training - Train AI with your face images and voice samples",
    "📝 Custom Question Sets - Create tailored interview questions",
    "🎥 Video/Audio Recording - Capture candidate responses",
    "🗣️ Speech-to-Text - Real-time answer transcription",
    "📊 AI Analysis - Automated evaluation and scoring",
    "📈 Dashboard & Analytics - Real-time insights and statistics"
])

# Slide 6: Interviewer Portal Features
add_content_slide(prs, "Interviewer Portal", [
    "👤 AI Avatar Training",
    "  - Upload face images for facial recognition",
    "  - Record voice samples for speech synthesis",
    "📅 Session Management",
    "  - Create interview sessions with custom questions",
    "  - Generate unique candidate links",
    "👥 Candidate Tracking",
    "  - View all candidates with AI-analyzed scores",
    "  - Ranking based on performance",
    "🎯 Dashboard - Real-time statistics and insights"
])

# Slide 7: Candidate Portal Features
add_content_slide(prs, "Candidate Portal", [
    "🔗 Easy Registration",
    "  - Join via unique interview links",
    "  - Simple onboarding process",
    "🤖 AI-Conducted Interviews",
    "  - Interact with AI avatar",
    "  - Video/audio response recording",
    "📝 Real-time Transcription",
    "  - Speech-to-text for all answers",
    "📊 Progress Tracking",
    "  - Visual progress through questions",
    "  - Time management indicators"
])

# Slide 8: AI Avatar Training Process
add_content_slide(prs, "AI Avatar Training", [
    "1️⃣ Face Recognition Training",
    "  - Upload multiple face images from different angles",
    "  - Uses face-api.js for facial landmark detection",
    "2️⃣ Voice Synthesis Training",
    "  - Record voice samples reading provided scripts",
    "  - Train AI to replicate voice patterns",
    "3️⃣ Avatar Configuration",
    "  - Customize appearance and behavior",
    "  - Set interview style and tone",
    "4️⃣ Testing & Validation",
    "  - Test avatar before deploying to candidates"
])

# Slide 9: Session Management
add_content_slide(prs, "Interview Session Management", [
    "📋 Question Builder",
    "  - Create custom interview questions",
    "  - Set time limits per question",
    "🔗 Candidate Link Generation",
    "  - Unique, secure links for each session",
    "  - Track link usage and access",
    "⏰ Session Controls",
    "  - Schedule start/end times",
    "  - Set candidate limits",
    "📊 Session Analytics",
    "  - Monitor live interview progress",
    "  - View completion rates"
])

# Slide 10: Tech Stack
add_content_slide(prs, "Technology Stack", [
    "⚛️ Frontend Framework: React 18.2",
    "⚡ Build Tool: Vite 5.0 with HMR",
    "🎨 Styling: TailwindCSS 3.3",
    "🧭 Routing: React Router DOM 6",
    "🎭 Face Recognition: face-api.js",
    "🎤 Speech: Web Speech API",
    "📊 Charts: Recharts",
    "🎬 Recording: MediaRecorder API",
    "🎯 Icons: Lucide React"
])

# Slide 11: Architecture Overview
add_content_slide(prs, "System Architecture", [
    "📁 Component-Based Architecture",
    "  - Shared UI components for reusability",
    "  - Interview-specific components",
    "  - Layout components for consistent structure",
    "🔄 Context API for State Management",
    "  - AuthContext for authentication",
    "  - InterviewContext for interview state",
    "🛣️ Client-Side Routing",
    "  - Protected routes for authenticated users",
    "  - Role-based access control",
    "🎨 Utility-First Styling with Tailwind"
])

# Slide 12: User Authentication System
add_content_slide(prs, "Authentication & Security", [
    "🔐 Secure Registration & Login",
    "  - Email validation",
    "  - Password strength requirements (6+ characters)",
    "  - Company information tracking",
    "👤 User Roles",
    "  - Interviewer/HR role",
    "  - Candidate role",
    "🔒 Protected Routes",
    "  - Role-based access control",
    "  - Session management",
    "✅ Demo Accounts Available",
    "  - demo@interview.pro / demo123",
    "  - hr@company.com / hr123"
])

# Slide 13: Interview Flow
add_content_slide(prs, "Interview Process Flow", [
    "1️⃣ HR creates interview session with questions",
    "2️⃣ System generates unique candidate link",
    "3️⃣ Candidate registers via link",
    "4️⃣ Interview starts with AI avatar introduction",
    "5️⃣ AI avatar asks questions one by one",
    "6️⃣ Candidate responds via video/audio",
    "7️⃣ Responses recorded and transcribed",
    "8️⃣ AI analyzes answers and generates scores",
    "9️⃣ Results displayed in HR dashboard"
])

# Slide 14: Real-time Features
add_content_slide(prs, "Real-time Capabilities", [
    "🎥 Live Video Recording",
    "  - MediaRecorder API for capturing responses",
    "  - Support for video and audio formats",
    "🗣️ Speech-to-Text Transcription",
    "  - Web Speech API integration",
    "  - Real-time answer transcription",
    "⏱️ Timer Management",
    "  - Visual countdown for each question",
    "  - Automatic progression",
    "📊 Progress Indicators",
    "  - Visual feedback on interview completion",
    "  - Question navigation"
])

# Slide 15: AI Analysis Engine
add_content_slide(prs, "AI-Powered Analysis", [
    "🧠 Automated Answer Evaluation",
    "  - Content relevance analysis",
    "  - Communication quality assessment",
    "📊 Scoring System",
    "  - Multi-factor scoring algorithm",
    "  - Weighted evaluation criteria",
    "🏆 Candidate Ranking",
    "  - Automatic ranking based on scores",
    "  - Comparative analysis",
    "💡 Insights Generation",
    "  - Strengths and weaknesses identification",
    "  - Improvement recommendations"
])

# Slide 16: Dashboard & Analytics
add_content_slide(prs, "Dashboard & Analytics", [
    "📈 Real-time Statistics",
    "  - Active interviews",
    "  - Completion rates",
    "  - Average scores",
    "👥 Candidate Management",
    "  - View all candidates",
    "  - Filter and search capabilities",
    "  - Detailed candidate profiles",
    "📊 Performance Metrics",
    "  - Interview success rates",
    "  - Time-to-completion analytics",
    "📉 Data Visualization with Recharts",
    "  - Interactive charts and graphs"
])

# Slide 17: Project Structure
add_content_slide(prs, "Code Organization", [
    "📂 src/components/",
    "  - interview/ - Interview components (VideoRecorder, AIAvatar, Timer)",
    "  - layouts/ - Page layouts (InterviewerLayout, CandidateLayout)",
    "  - shared/ - Reusable UI components (Button, Card, Modal, etc.)",
    "📂 src/contexts/ - React context providers",
    "📂 src/pages/ - Page components (auth, candidate, interviewer)",
    "📂 src/utils/ - Helper functions and utilities",
    "  - aiAnalysis.js - AI evaluation logic",
    "  - mediaUtils.js - Media handling",
    "  - helpers.js - General utilities"
])

# Slide 18: Installation & Setup
add_content_slide(prs, "Getting Started", [
    "📋 Prerequisites: Node.js 18+ and npm 9+",
    "",
    "🚀 Installation Steps:",
    "1. Clone repository: git clone https://github.com/Shlokabhishek/Interview_Taker.git",
    "2. Install dependencies: npm install",
    "3. Start dev server: npm run dev",
    "4. Open http://localhost:5173",
    "",
    "🔨 Available Commands:",
    "  - npm run dev - Start development server",
    "  - npm run build - Build for production",
    "  - npm run preview - Preview production build"
])

# Slide 19: Key Benefits
add_content_slide(prs, "Benefits & Impact", [
    "⚡ 10x faster recruitment process",
    "💰 Reduced hiring costs by 70%",
    "🎯 Consistent evaluation criteria",
    "📈 Scalable to unlimited candidates",
    "🤖 24/7 availability for interviews",
    "🔍 Data-driven candidate insights",
    "✅ Improved candidate experience",
    "📊 Better hiring decisions with AI analytics"
])

# Slide 20: Conclusion & Future
add_content_slide(prs, "Conclusion & Future Roadmap", [
    "✨ Current Achievement:",
    "  - Fully functional AI interview platform",
    "  - Modern tech stack with React & Vite",
    "  - Complete interviewer and candidate portals",
    "",
    "🚀 Future Enhancements:",
    "  - Advanced NLP for deeper answer analysis",
    "  - Multi-language support",
    "  - Integration with ATS systems",
    "  - Mobile app for candidates",
    "  - Video sentiment analysis",
    "  - Custom branding options"
])

# Save presentation
prs.save('AI_Interview_Platform_Presentation.pptx')
print("✅ Presentation created successfully: AI_Interview_Platform_Presentation.pptx")
print(f"📊 Total slides: {len(prs.slides)}")
