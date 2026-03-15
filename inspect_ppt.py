from pptx import Presentation

prs = Presentation("PPT_Template.pptx")

with open("ppt_template_dump.txt", "w", encoding="utf-8") as f:
    f.write(f"slides {len(prs.slides)}\n")
    for i, slide in enumerate(prs.slides, 1):
        f.write(f"--- Slide {i} ---\n")
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip().replace("\n", " | "))
        f.write(("\n".join(texts) if texts else "[no text]") + "\n")
