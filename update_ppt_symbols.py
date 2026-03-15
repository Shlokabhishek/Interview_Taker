from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt

INPUT_FILE = "AI_Viewers_Hackathon_Presentation_v2.pptx"
OUTPUT_FILE = "AI_Viewers_Hackathon_Presentation_v2.pptx"
FALLBACK_OUTPUT_FILE = "AI_Viewers_Hackathon_Presentation_v2_symbols.pptx"


def _rgb(color: tuple[int, int, int]) -> RGBColor:
    return RGBColor(color[0], color[1], color[2])


def delete_shapes(slide, predicate) -> int:
    removed = 0
    for shape in list(slide.shapes):
        try:
            if predicate(shape):
                el = shape._element
                el.getparent().remove(el)
                removed += 1
        except Exception:
            # Best-effort deletion; avoid breaking the update if a single shape is odd.
            continue
    return removed


def add_labeled_shape(
    slide,
    shape_type,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    *,
    fill: tuple[int, int, int],
    font_color: tuple[int, int, int] = (255, 255, 255),
    font_size: int = 14,
    bold: bool = True,
    line: tuple[int, int, int] | None = (255, 255, 255),
):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line)

    tf = shape.text_frame
    tf.clear()
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _rgb(font_color)
    return shape


def add_label(slide, text: str, left: int, top: int, width: int, height: int, *, font_size: int = 16):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(34, 34, 34)
    return box


def add_right_arrow(slide, left: int, top: int, width: int, height: int, *, fill=(200, 200, 200)):
    arrow = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = _rgb(fill)
    arrow.line.color.rgb = _rgb(fill)
    return arrow


def update_slide_4(slide):
    # Remove the existing "rounded rectangle + '->' text" flow chart (keep headings/summary).
    flow_texts = {"HR Login", "Create Session", "Share Link", "AI Interview", "Score & Rank", "->"}

    def is_old_flow_shape(shape) -> bool:
        if not getattr(shape, "has_text_frame", False):
            return False
        text = (shape.text_frame.text or "").strip()
        if text not in flow_texts:
            return False
        # Restrict to the flowchart band so we don't delete other arrows/text.
        return 1600000 <= shape.top <= 2600000

    delete_shapes(slide, is_old_flow_shape)

    # Build a proper flow chart using standard symbols.
    top = 1850000
    height = 700000
    arrow_h = 320000
    arrow_top = top + (height - arrow_h) // 2

    start_end_w = 900000
    proc_w = 1350000
    arrow_w = 260000
    gap = 60000

    x = 650000

    add_labeled_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
        "Start",
        x,
        top,
        start_end_w,
        height,
        fill=(75, 85, 99),
        line=None,
    )
    x += start_end_w + gap
    add_right_arrow(slide, x, arrow_top, arrow_w, arrow_h)
    x += arrow_w + gap

    steps = [
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "HR Login", (34, 87, 122)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Create Session", (46, 125, 50)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Share Link", (255, 143, 0)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PREDEFINED_PROCESS, "AI Interview", (142, 36, 170)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Score & Rank", (198, 40, 40)),
    ]

    for idx, (shape_type, label, color) in enumerate(steps):
        add_labeled_shape(
            slide,
            shape_type,
            label,
            x,
            top,
            proc_w,
            height,
            fill=color,
            line=None,
            font_size=14,
        )
        x += proc_w + gap
        add_right_arrow(slide, x, arrow_top, arrow_w, arrow_h)
        x += arrow_w + gap

    add_labeled_shape(
        slide,
        MSO_AUTO_SHAPE_TYPE.FLOWCHART_TERMINATOR,
        "End",
        x,
        top,
        start_end_w,
        height,
        fill=(75, 85, 99),
        line=None,
    )


def update_slide_8(slide):
    # Remove old phase list + pipeline text (replace with diagrams).
    def is_old_roadmap_text(shape) -> bool:
        if not getattr(shape, "has_text_frame", False):
            return False
        text = (shape.text_frame.text or "").strip()
        return text.startswith("Phase 1:") or text.startswith("Pipeline:")

    delete_shapes(slide, is_old_roadmap_text)

    add_label(slide, "Roadmap", left=1180000, top=1600000, width=2500000, height=260000, font_size=16)

    # Roadmap timeline (flowchart process symbols connected by arrows).
    top = 1880000
    height = 760000
    box_w = 1600000
    arrow_w = 220000
    arrow_h = 260000
    arrow_top = top + (height - arrow_h) // 2
    gap = 50000

    x = 1180000
    phases = [
        ("Phase 1\nAuth + Sessions", (34, 87, 122)),
        ("Phase 2\nAvatar Training", (46, 125, 50)),
        ("Phase 3\nRecording + STT", (255, 143, 0)),
        ("Phase 4\nScoring + Ranking", (142, 36, 170)),
        ("Phase 5\nDashboard + Deploy", (198, 40, 40)),
    ]

    for idx, (label, color) in enumerate(phases):
        add_labeled_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS,
            label,
            x,
            top,
            box_w,
            height,
            fill=color,
            line=None,
            font_size=12,
        )
        x += box_w + gap
        if idx < len(phases) - 1:
            add_right_arrow(slide, x, arrow_top, arrow_w, arrow_h, fill=(190, 190, 190))
            x += arrow_w + gap

    add_label(
        slide, "Algorithm Pipeline", left=1180000, top=2850000, width=3200000, height=260000, font_size=16
    )

    # Pipeline (I/O -> processes -> document) + data store.
    pipe_top = 3150000
    pipe_h = 760000
    node_w = 1300000
    node_h = pipe_h
    arrow_w = 200000
    arrow_h = 240000
    arrow_top = pipe_top + (node_h - arrow_h) // 2
    gap = 50000

    x = 1180000
    nodes = [
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_DATA, "Input\nAnswer", (0, 128, 128)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Transcription\n(STT)", (55, 65, 81)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Feature\nExtraction", (55, 65, 81)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Score\nCalculation", (55, 65, 81)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_PROCESS, "Ranking", (55, 65, 81)),
        (MSO_AUTO_SHAPE_TYPE.FLOWCHART_DOCUMENT, "Dashboard\nOutput", (34, 87, 122)),
    ]

    anchor_for_db_x = None
    for idx, (shape_type, label, color) in enumerate(nodes):
        add_labeled_shape(
            slide,
            shape_type,
            label,
            x,
            pipe_top,
            node_w,
            node_h,
            fill=color,
            line=None,
            font_size=12,
        )
        if idx == 2:
            anchor_for_db_x = x
        x += node_w + gap
        if idx < len(nodes) - 1:
            add_right_arrow(slide, x, arrow_top, arrow_w, arrow_h, fill=(190, 190, 190))
            x += arrow_w + gap

    # Data store under the pipeline
    if anchor_for_db_x is not None:
        db_left = anchor_for_db_x + 250000
        db_top = 4100000
        db_w = 1300000
        db_h = 550000
        add_labeled_shape(
            slide,
            MSO_AUTO_SHAPE_TYPE.FLOWCHART_STORED_DATA,
            "Responses\nDB",
            db_left,
            db_top,
            db_w,
            db_h,
            fill=(230, 233, 238),
            font_color=(34, 34, 34),
            bold=True,
            font_size=12,
            line=(200, 205, 215),
        )

        # Simple connectors down from Feature Extraction and Score Calculation
        line1 = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            anchor_for_db_x + node_w // 2,
            pipe_top + node_h,
            db_left + db_w // 2,
            db_top,
        )
        line1.line.color.rgb = _rgb((200, 205, 215))
        line1.line.width = Pt(2)


def main() -> int:
    prs = Presentation(INPUT_FILE)

    # Slide indexes are 0-based. We target:
    # - Slide 4 (index 3): interview flow chart
    # - Slide 8 (index 7): roadmap + pipeline
    if len(prs.slides) < 8:
        raise SystemExit(f"Unexpected deck size: {len(prs.slides)} slides (need at least 8)")

    update_slide_4(prs.slides[3])
    update_slide_8(prs.slides[7])

    try:
        prs.save(OUTPUT_FILE)
        print(f"Updated {OUTPUT_FILE}")
    except PermissionError:
        fallback = Path(FALLBACK_OUTPUT_FILE)
        prs.save(fallback)
        print(f"Updated {fallback}")

    return 0


if __name__ == "__main__":
    raise SystemExit(main())

